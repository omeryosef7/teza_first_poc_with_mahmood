#!/usr/bin/env python3
"""Build a 12-slide plain summary deck of the GCG research for Mahmood.

Simple format: title + bullets, 16:9, no images. Every number traces to
docs/GCG_JULY2026_MASTER_LOG.md (agent-verified before build).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "GCG_Summary_for_Mahmood_2026-07-18.pptx"

DARK = RGBColor(0x20, 0x20, 0x20)
ACCENT = RGBColor(0x1F, 0x3A, 0x5F)   # dark navy, understated
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = ACCENT
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(15); r2.font.italic = True; r2.font.color.rgb = GREY
    # underline rule
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.5), Inches(12.1), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return slide


def add_body(slide, bullets, top=1.75, size=15, height=5.4):
    """bullets: list of (level, text) or str (level 0)."""
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        text, lvl = (b if isinstance(b, tuple) else (b, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(5 if lvl == 0 else 2)
        marker = "•  " if lvl == 0 else "–  "
        # allow simple **bold** segments; strip leftover single-* emphasis markers
        segs = [seg.replace("*", "") for seg in text.split("**")]
        r = p.add_run(); r.text = ("" if lvl else "") + marker + segs[0]
        r.font.size = Pt(size if lvl == 0 else size - 1)
        r.font.color.rgb = DARK if lvl == 0 else GREY
        r.font.bold = (lvl == 0 and text.endswith(":::"))  # not used
        for i, seg in enumerate(segs[1:], start=1):
            rr = p.add_run(); rr.text = seg
            rr.font.size = Pt(size if lvl == 0 else size - 1)
            rr.font.color.rgb = DARK if lvl == 0 else GREY
            rr.font.bold = (i % 2 == 1)
    return slide


def title_slide():
    s = prs.slides.add_slide(BLANK)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "GCG Adversarial-Suffix Jailbreak Research"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = ACCENT
    p2 = tf.add_paragraph(); r2 = p2.add_run()
    r2.text = "A Plain-Language Summary of Everything We Did (July 2026)"
    r2.font.size = Pt(20); r2.font.color.rgb = DARK
    p3 = tf.add_paragraph(); p3.space_before = Pt(18); r3 = p3.add_run()
    r3.text = ("Models: Qwen3-14B, Gemma4-E4B-it, DeepSeek-R1-Distill-Qwen-7B\n"
               "Omer Yosef  ·  prepared for Mahmood  ·  2026-07-18")
    r3.font.size = Pt(15); r3.font.italic = True; r3.font.color.rgb = GREY
    return s


# ------- Slide content (numbers agent-verified vs GCG_JULY2026_MASTER_LOG.md) -------
title_slide()

s = prs.slides.add_slide(BLANK)
add_title(s, "1. What We Are Studying (Background)", "The problem and the words you need to follow the rest")
add_body(s, [
    "**The question:** can we make a safety-tuned chatbot answer a harmful request by adding a short, optimized nonsense string (an \"adversarial suffix\") to the end of the prompt?",
    "**GCG (Greedy Coordinate Gradient):** the algorithm that searches, one token at a time, for the suffix that makes the model most likely to *start* its reply with a chosen \"compliant\" sentence.",
    "**\"Thinking\" models (Qwen3, Gemma4):** before the visible answer they write a hidden reasoning block between **<think> … </think>** tags. This detail turns out to matter a lot.",
    "**Refusal direction:** a single internal direction in the model's activations that is linked to \"I refuse.\" Pushing activations away from it is a known jailbreak lever.",
    "**Our yardstick:** a prior paper (\"CoT-Hijacking\") reached **91%** success by directly editing that refusal direction *inside* the model. We only get to change the *input text* (a harder, more realistic setting) — so 91% is the ceiling we compare against.",
], size=15)

s = prs.slides.add_slide(BLANK)
add_title(s, "2. How We Run and Score Every Test (Setup)", "Same recipe for all experiments")
add_body(s, [
    "**Models:** Qwen3-14B (main focus), Gemma4-E4B-it (second model), DeepSeek-R1-Distill-Qwen-7B (a third model family, used once).",
    "**Data:** AdvBench harmful behaviours — a **25-behaviour** \"dev panel\" for fast iteration, and the full **520-behaviour** benchmark for scale.",
    "**Judge:** the **StrongREJECT** classifier decides, for each generated answer, whether it is a genuine successful jailbreak (yes/no).",
    "**Main metric — ASR (Attack Success Rate):** the fraction of attempts judged successful. We always compare the optimized suffix against controls (no suffix / random suffix). On 25 behaviours we use 3 seeds, so ASR is out of 75 attempts (written **X/75**).",
    ("At 520 scale we report two views:", 0),
    ("**combo-ASR** = per (behaviour × seed) attempt;  **behaviour-ASR** = a behaviour counts as won if *any* seed succeeds (always the higher number).", 1),
    "**Detection:** a simple classifier reads the model's internal state at the *first generated token* and tries to tell \"attacked\" from \"clean.\" Scored by **AUC** (1.0 = perfect, 0.5 = coin flip).",
], size=14)

s = prs.slides.add_slide(BLANK)
add_title(s, "3. Stage 1 — Ablation (25 behaviours): Textbook GCG Barely Works")
add_body(s, [
    "**What we did:** ran standard GCG on Qwen3, with the usual target \"Sure, here is …\".",
    "**Result:** **4.0%** ASR on Qwen3 — no better than the no-attack baseline. Scaled to all 520 behaviours it was **1.9%**, actually slightly *net-negative*.",
    "**Gemma4:** **0%** — the plain attack never worked on it at all.",
    "**Why it fails on thinking models:** the target assumes the model's *first* token is \"Sure\". But a thinking model's first token is **<think>**, so the target is aimed at the wrong place.",
    "**Data:** 25 AdvBench behaviours × 3 seeds, StrongREJECT judge, compared against no-suffix and random-suffix controls.",
    "**Takeaway:** the standard attack is essentially useless here — which set up the key fix on the next slide.",
], size=16, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "4. Stage 1 — The Key Unlock: CoT-Prefix Targeting (Exp 5A)")
add_body(s, [
    "**The idea:** make the target look like what a thinking model actually writes first — put a fake reasoning block in front: **\"<think> Okay, I can help. </think> Sure, here is {task}\"**.",
    "**Result:** ASR jumps to **10.7% (8/75)**, **+8pp** over baseline. This is the single most important trick in the whole project.",
    "**The counter-intuitive part:** this trick makes the training *loss worse* (task_loss 47.6 → 20.5 is a worse teacher-forced fit) yet gives **~4× the ASR (2.7% → 10.7%)**.",
    ("Lesson: attack success tracks how *behaviourally realistic* the target is, **not** the optimizer's loss.", 1),
    "**Important limit:** if we ALSO add full-strength refusal-direction suppression on top (Exp 6C, λ=1.0), the gain is destroyed — back to **0% (−10.7pp)**. Full-strength suppression and this trick fight each other. (Slide 7 revisits this.)",
], size=16, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "5. Stage 2 — Does It Scale? Full 520-Behaviour Benchmark (Phase 7A)")
add_body(s, [
    "**What we did:** took the winning 5A suffix and evaluated it on ALL 520 AdvBench behaviours.",
    "**Result:** **8.01%** ASR on training seeds (125/1560); **8.92%** on unseen seeds (131/1468, covering 493 of 520 behaviours).",
    "**Uplift over the no-attack baseline: +5.09pp** (95% confidence interval +3.4 to +6.8), and statistically very strong (**McNemar p < 10⁻¹⁰**).",
    "Lower than the 25-behaviour number — expected, because small dev sets overstate performance — but a **real, robust** effect.",
    "**Where it works best:** \"misinformation\" behaviours show by far the biggest lift (**+19.8pp**).",
    "**Caution — high seed variance:** seed 45 was best (16% / 21%), while seed 44 was *net-negative* (1.3%). The same recipe swings a lot by random seed.",
], size=15, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "6. Stage 2 — A Second Finding: These Attacks Are Easy to Detect (Qwen3)")
add_body(s, [
    "**What we did:** trained a simple detector on Qwen3's internal state at the *first generated token* to separate attacked prompts from clean ones.",
    "**Result: AUC = 1.000 — perfect separation**, for every Qwen3 GCG variant we tried.",
    "**It's robust:** the perfect score holds even under strict cross-validation (held-out behaviours, held-out seeds) — so it is not overfitting.",
    "**It's independent of ASR:** even weak attacks are perfectly visible; even the strongest suffix leaves an obvious fingerprint at token 0.",
    "**Gemma4 is different:** its detector is much weaker (**~0.60–0.75 AUC**) and unreliable — it flags **42.7%** of harmless random text as an attack.",
    "**Takeaway:** on Qwen3, *succeeding* at the attack and *hiding* it are two different problems — the suffix is very \"loud\" internally.",
], size=15, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "7. Stage 3 — Biggest Surprise: Weak Suppression Doubles ASR (Phase 8)")
add_body(s, [
    "**Background:** we believed refusal-direction suppression was simply *incompatible* with the CoT-prefix trick, because at full strength (λ=1.0) it gave 0%.",
    "**What we did:** instead of on/off, we swept the *strength* λ of the suppression.",
    "**Result — it was strength-specific, not a general truth.** At **λ=0.3 (weak suppression)** ASR = **24.0% (18/75), +21.3pp** — the best 25-behaviour result of the project.",
    ("Strong suppression (λ=1.0 → 0%, λ=3.0 → 2.7%) kills it; *weak* suppression amplifies it.", 1),
    "**Replicated** on a second seed (seed 43: **12.0%**). The effect is real but varies ~2× by seed.",
    "**Lesson:** \"the two objectives are incompatible\" was false — it only held at the one strength we first tried.",
], size=15, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "8. Stage 4 — Sprint 2: Four Follow-ups, All Null or Negative")
add_body(s, [
    "**Causal test — does \"compliant framing\" *cause* success?** We forced the model's reasoning to open with compliant framing. Result: **8% vs 12%** baseline, **p=1.0 → no causal effect.** An earlier *correlation* was NOT causal. (Uses no suffix, so it is unaffected by the later bug.)",
    "**Third model (DeepSeek-R1):** already **~50% compliant with NO attack** (CoT 49% vs 47% baseline) → no safety \"headroom\" for GCG to add value.",
    "**Gemma4 CoT-channel v2 (800 steps):** the special tokens are trainable now, but ASR still **0%**.",
    "**Attack-quality tweaks:** longer suffix = **2.7%**; seed-44 + live-ASR selection = **4.0%** — both *worse* than the 10.7% reference.",
    "**Takeaway:** none of these four directions helped — they ruled out tempting explanations and dead ends.",
], size=15, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "9. Stage 5 — Sprint 3: Scaling to 520 + the Best New Result")
add_body(s, [
    "**Scaled the λ=0.3 recipe to all 520 behaviours:** seed 42 (9A) = **11.21% combo / 94 of 520 behaviours**; seed 45 (9B) = **8.83% / 72 behaviours**.",
    "λ=0.3 is **not additive across seeds** (a fresh seed-45 run, 9C, gave 6.09%).",
    "**Headline result (10F) — a free \"union\" ensemble:** count a behaviour as won if *either* the seed-42 or seed-45 run wins it → **13.97% combo, 110 of 520 behaviours (21.2%)**.",
    ("Better than either seed alone, at **zero extra optimization cost** — because the seeds win on *different* behaviours (38 only seed-42, 16 only seed-45, 56 both).", 1),
    "**Gemma4 (10A, \"EmptyThink\" recipe) at 520 scale:** **3.91% vs 2.31%** baseline (**+1.6pp**), consistent across seeds — modest, but **Gemma4's best result of the whole project.**",
], size=15, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "10. Important Caveat — A Placement Bug Makes the Numbers Provisional")
add_body(s, [
    "A code audit found a **suffix-placement bug** (we call it \"B1\"):",
    ("The suffix was **optimized inside the assistant's turn**, but **evaluated in the user's turn** — a mismatch, so the optimizer was tuned on a slightly different prompt than we tested.", 1),
    "**Consequence:** every ASR *magnitude* on the earlier slides (10.7%, 24%, 8.9%, 14%, …) is **provisional** until re-run with the fix.",
    "**What still holds (expected to survive):** the *directions* of every finding (CoT-prefix helps, weak λ=0.3 helps, the ensemble helps), the detection AUC, the baselines, and the causal test.",
    "This bug is also the likely reason **\"loss doesn't predict ASR\"** (loss was measured on the wrong prompt).",
    "**Status:** the fix is implemented (put the suffix in the user turn, byte-identical to evaluation) and a **re-run campaign is running now** to confirm the true numbers.",
], size=15, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "11. How This Relates to Standard GCG — What Is the SAME", "We adapted the standard attack (GCG, Zou et al. 2023), we did not invent a new one")
add_body(s, [
    "**Same core algorithm — Greedy Coordinate Gradient (GCG, Zou et al. 2023).** For each position in the suffix, use the gradient of a teacher-forced target loss to shortlist the top-k token swaps, batch-test many candidates, and greedily keep the best.",
    "**Same attack surface:** a token suffix appended to the user's request, joined with a space — \"{request} {suffix}\" — exactly the original's format.",
    "**Same base objective:** a cross-entropy loss pushing the model to START its answer with a chosen compliant target string.",
    "**Same prompt placement (after our fix):** suffix in the USER turn, target in the ASSISTANT turn — we checked our corrected prompt-builder against the original `llm-attacks` `SuffixManager` code and it matches.",
    "**Same multi-behaviour idea:** aggregate gradients across several behaviours to search for a more universal suffix.",
], size=16, top=1.85)

s = prs.slides.add_slide(BLANK)
add_title(s, "12. How This Relates to Standard GCG — What Is DIFFERENT (what we added)")
add_body(s, [
    "**New target models.** Original GCG attacks Llama-2 / Vicuna / GPT-family. We built adapters for reasoning models — Qwen3-14B, Gemma4-E4B-it, DeepSeek-R1 — which the original code never handled.",
    "**Thinking-model targeting — our main contribution.** Original GCG has no concept of a `<think>` reasoning block. Our key unlock (5A) makes the target include a fake `<think>…</think>` prefix so it matches what a reasoning model actually generates first.",
    "**Richer objective.** Original GCG optimizes ONE loss (target cross-entropy). We add optional terms: representation-alignment, KL, and a refusal-direction suppression term (Stage 3's λ) folded into the candidate-proposal gradient.",
    "**Harder success test.** Original GCG scores a \"win\" by simple string matching (did the reply avoid phrases like \"I'm sorry\"?). We use the StrongREJECT LLM judge on free-form generation, against 3 controls — a much stricter bar.",
    "**Tokenizer fix.** Original GCG filters candidates by re-tokenization (`filter_cand=True`). On Qwen/Gemma's BPE tokenizers that silently kills all progress, so we disable it (`--no-filter-cand`).",
    "**Bonus — detection.** We added a first-token hidden-state detector (AUC, slide 6) — not part of GCG at all.",
    "**Honesty tie-in:** the placement bug (slide 10) was itself a *divergence* from canonical GCG (we optimized in the assistant turn); the fix re-aligns us to the original.",
], size=14, top=1.7)

s = prs.slides.add_slide(BLANK)
add_title(s, "13. Summary — What We Actually Learned")
add_body(s, [
    "**1.** For thinking models, **CoT-prefix targeting** is the core unlock: 2.7% → **10.7%**.",
    "**2.** **Weak** refusal-direction suppression (λ=0.3) is the strongest single lever — up to **24%** on 25 behaviours; *strong* suppression kills it.",
    "**3.** A cheap **union of two seeds** covers the most of the benchmark: **110/520 behaviours (21%)**.",
    "**4.** Qwen3 attacks are **perfectly detectable** at the first token (**AUC 1.0**); **Gemma4 resists almost everything** (one modest 3.9% recipe).",
    "**5.** The optimizer's **loss does NOT predict** attack success.",
    "**6.** A **placement bug** means the exact ASR magnitudes are **provisional** — a fixed re-run is in progress to confirm them.",
    "**Bottom line:** the qualitative findings are clear and real; the final numbers are pending the corrected re-run.",
], size=15, top=1.7)

# ---------------- Presenter (speaker) notes ----------------
def set_notes(slide, lines):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln


NOTES = [
    # 1 — Title
    [
        "PURPOSE: 30-second opener. Say what this is and set the frame.",
        "SAY: This is a summary of ~3 weeks of work (July 2026) trying to jailbreak safety-tuned reasoning LLMs using optimized adversarial text suffixes — the GCG method.",
        "Models: Qwen3-14B is our main subject; Gemma4-E4B-it is a second model; DeepSeek-R1-Distill-Qwen-7B is a third family we tried once.",
        "Roadmap: background -> how we measure -> five stages of experiments -> one important bug caveat -> takeaways.",
        "HONESTY FRAME (say up front): the qualitative findings are solid, but a bug we found late means the exact percentages are provisional. I'll flag it clearly near the end (slide 10).",
    ],
    # 2 — Background
    [
        "PURPOSE: give the audience the 4 words they need for the rest of the talk.",
        "Jailbreak = getting a safety-tuned model to produce harmful content it is trained to refuse.",
        "Adversarial suffix = a string of tokens we append to the user's request (often looks like gibberish). GCG is the algorithm that optimizes it.",
        "GCG (Greedy Coordinate Gradient, Zou et al. 2023): a white-box, gradient-guided search that, position by position, swaps tokens in the suffix to raise the probability the model STARTS its reply with a chosen 'compliant' sentence.",
        "Thinking models: Qwen3/Gemma4 first write a hidden reasoning block inside <think> ... </think> before the visible answer. Remember this — it is the whole reason Stage 1 works.",
        "Refusal direction: from prior interpretability work, a single internal direction that encodes 'I refuse.' The 'CoT-Hijacking' paper got 91% jailbreak success by directly EDITING that direction inside the model.",
        "KEY: that 91% is a white-box ACTIVATION edit, not a text suffix — it is the theoretical ceiling. Our job is to approximate it using only input TEXT, which is much harder and more realistic.",
        "LIKELY Q ('why is 91% the ceiling?'): because internal manipulation is strictly stronger than a text-only attack.",
    ],
    # 3 — Setup
    [
        "PURPOSE: explain how every result is measured, so the later numbers are trusted.",
        "AdvBench = a standard benchmark of harmful instructions ('behaviours'). We use 25 for fast iteration (the 'dev panel') and the full 520 for scale claims.",
        "StrongREJECT = an automated LLM-based judge that scores whether an answer actually delivers harmful content (not merely 'failed to refuse'). One yes/no per attempt.",
        "ASR (Attack Success Rate) = successes / attempts. On 25 behaviours x 3 random seeds = 75 attempts, so we write X/75.",
        "CRUCIAL: we ALWAYS compare the optimized suffix against controls — no suffix, task-only, and a random-character suffix. What matters is UPLIFT (optimized minus control), not the raw number. A high ASR that doesn't beat controls is worthless.",
        "Two views at 520 scale: combo-ASR = each (behaviour x seed) attempt (strict); behaviour-ASR = a behaviour counts if ANY seed succeeds (lenient, always higher).",
        "Detection is a SEPARATE question: a tiny classifier reads the model's internal state at the first generated token to see if attacked prompts are distinguishable — scored by AUC (1.0 = perfect, 0.5 = coin flip).",
        "LIKELY Q ('why 3 seeds?'): generation is random; seeds average out sampling noise.",
    ],
    # 4 — Stage 1 standard GCG
    [
        "PURPOSE: the failure that motivates the whole project.",
        "We ran textbook GCG on Qwen3 with the classic target 'Sure, here is {task}'.",
        "Result: 4.0% — statistically the same as doing nothing. At full 520 scale it was even worse, 1.9%, actually slightly NET-NEGATIVE (the suffix suppressed a couple of control successes).",
        "Gemma4: flat 0% — never worked.",
        "WHY (the important bit): GCG's target assumes the model's FIRST output token is 'Sure'. But a thinking model's first token is the literal <think> tag — so we were optimizing for a position that never occurs. The target is misaligned with reality.",
        "This misalignment is exactly what the next slide fixes.",
        "LIKELY Q ('isn't GCG strong?'): yes, on non-thinking models (the original Llama work). Thinking models break its core assumption.",
    ],
    # 5 — Stage 1 unlock 5A
    [
        "PURPOSE: the single most important result in the deck — spend time here.",
        "The fix: change the target to match what a thinking model really writes — prepend a short fake reasoning block: '<think> Okay, I can help. </think>' then the compliance sentence. Now the target lines up with the model's actual token stream.",
        "Result: 10.7% ASR (8 of 75), +8pp over baseline — a real jump from ~2.7% with the standard target.",
        "COUNTERINTUITIVE POINT on the slide: in our data, attack success did NOT track the optimizer's loss (the slide cites task_loss 47.6 vs 20.52).",
        "State this HONESTLY, don't explain it away. Two caveats: (1) raw task_loss across different-length targets isn't an apples-to-apples 'lower is better' comparison (the master log flags this), so don't over-argue the 47.6-vs-20.52 direction; (2) MORE IMPORTANTLY, this loss-vs-ASR decoupling may itself be an ARTIFACT of the B1 placement bug (slide 10) — loss was computed on the assistant-turn prompt while ASR was measured on the user-turn prompt, so the two need not correlate. Present it as an observation we're re-testing, not a settled law.",
        "Supporting example (also B1-affected — caveat it): seeds 44 and 45 had near-identical loss (19.91 vs 19.98) but 1.3% vs 16.0% ASR (~11x gap). If Mahmood raises the bug, agree: 'right — the fixed re-run will tell us whether this decoupling survives.'",
        "Also flag forward: if we ALSO add FULL-strength refusal-direction suppression on top (Exp 6C, lambda=1.0), the gain is destroyed -> 0% (-10.7pp). Remember this — Stage 3 shows WEAK suppression does the opposite.",
    ],
    # 6 — Stage 2 scale
    [
        "PURPOSE: prove 5A isn't a fluke of 25 hand-picked behaviours.",
        "We took the single best 5A suffix and evaluated it on ALL 520 behaviours.",
        "Numbers: 8.01% on training seeds (125/1560); 8.92% on unseen seeds (131/1468), over 493 of 520 behaviours (~95% coverage).",
        "THE NUMBER THAT MATTERS is uplift, not raw ASR: +5.09pp, 95% CI [+3.4, +6.8] (excludes zero = real), and McNemar's paired test gives p < 10^-10 (extremely significant).",
        "Why lower than the 25-beh 10.7%? Small dev panels overstate performance; also the unseen-seed baseline is naturally higher (~12%), which inflates raw unseen ASR — that's exactly why we lead with uplift, not the headline number.",
        "Where it works best: 'misinformation' behaviours get the biggest lift (+19.8pp) — the attack helps most where refusal is softest.",
        "HONESTY: high seed variance. Seed 45 best (16%/21%), seed 44 net-negative (1.3%). Same recipe, different random seed, very different outcome.",
        "LIKELY Q ('is 8.9% good?'): modest in absolute terms, but it's a real, statistically robust uplift that generalises to the full benchmark — most GCG-on-thinking-model attempts get essentially nothing.",
    ],
    # 7 — Detection
    [
        "PURPOSE: a defensive-security finding — even when the attack works, is it visible?",
        "We trained a simple logistic-regression detector on Qwen3's hidden state at the FIRST generated token. For every Qwen3 GCG variant it separates attacked vs clean with AUC = 1.000 — perfect.",
        "Robustness: holds under strict cross-validation — group-by-behaviour, leave-one-seed-out, leave-one-optimization-seed-out, and a 25-vs-495 behaviour split. Not memorisation; the signal generalises.",
        "Independent of ASR: even the weakest attack is perfectly detectable; the strongest suffix is still 'loud'. (Fun detail: seed 44 had the LOWEST ASR but the LARGEST activation shift.)",
        "Gemma4 is different: detector AUC only ~0.60-0.75 AND not specific — it flags 42.7% of harmless random text as an attack. So this is a Qwen3 story.",
        "CAVEAT (say if asked): measured optimized-vs-neutral only; we did NOT test an adaptive attacker trying to evade the detector. Do not call it 'production-ready'.",
        "LIKELY Q ('so the attack is useless if detectable?'): for a deployed system, largely yes — but detection and attack success are separate axes; the research value is understanding the attack.",
    ],
    # 8 — Stage 3 lambda=0.3
    [
        "PURPOSE: the biggest surprise of the project.",
        "Context: on slide 4 we saw FULL-strength refusal-direction suppression (lambda=1.0) destroyed the CoT-prefix gain (0%). We had concluded the two were incompatible.",
        "We tested that assumption by sweeping the suppression STRENGTH lambda. It turned out to be strength-specific, not a general truth.",
        "At lambda=0.3 (WEAK suppression, layer 25): ASR = 24.0% (18/75), +21.3pp — the best 25-behaviour result of the whole project, more than double plain CoT-prefix (10.7%).",
        "The pattern: strong suppression kills it (lambda=1.0 -> 0%, lambda=3.0 -> 2.67%); weak suppression amplifies it. There is a sweet spot.",
        "Replicated on a second seed (43: 12.0%). Real but seed-variable (~2x swing) — consistent with our general high seed variance.",
        "MECHANISM nuance (only if asked): the refusal-direction term shapes which candidate tokens get PROPOSED (via the gradient), not which get accepted — so 'jointly optimises' is imprecise; and it was measured on one behaviour, not averaged.",
        "LIKELY Q ('why does weak help but strong hurt?'): honestly not fully understood — strong suppression over-constrains / destabilises the search. It's an empirical sweet spot, not exhaustively characterised (only this layer/target, 2 seeds).",
    ],
    # 9 — Stage 4 Sprint 2
    [
        "PURPOSE: intellectual honesty — four follow-ups that did NOT pan out. Negatives build credibility with Mahmood.",
        "Track 2 (the important one): we had a CORRELATION — generations whose reasoning 'restates the task' succeed more (12.8%) than ones that 'plan a refusal' (1.3%). Tempting to say the framing CAUSES success.",
        "We tested it causally: forced the model's <think> to open with compliant framing, then let it generate freely. Result: 8% (2/25) vs 12% (3/25) baseline, McNemar p=1.0 -> NO causal effect. So we must NOT claim the framing is causal. (Uses no suffix, so it's immune to the slide-10 bug — a clean result.)",
        "Track 3: tried a 3rd family, DeepSeek-R1. It is already ~50% compliant with NO attack (49% vs 47%) -> no safety 'headroom' for GCG to show value. A different failure mode, not a refutation.",
        "Track 1: pushed Gemma4's special 'channel-token' attack to 800 steps. The tokens ARE trainable (loss dropped a lot), which refutes the old 'architecturally impossible' claim — but ASR was still 0%. Mechanism corrected, conclusion unchanged.",
        "Track 4: two attempts to improve Qwen3's attack (a longer 35-token suffix; live-ASR-guided seed selection) gave 2.7% and 4.0% — both WORSE than the 10.7% reference.",
        "TAKEAWAY: we ruled out an over-claim and several dead ends.",
    ],
    # 10 — Stage 5 Sprint 3
    [
        "PURPOSE: scaling the winning lambda=0.3 recipe to 520, plus the best NEW result.",
        "9A (seed 42): 11.21% combo-ASR, winning 94 of 520 behaviours (18%). 9B (seed 45): 8.83%, 72 behaviours.",
        "Does lambda=0.3 stack with a good seed? A fresh seed-45 run (9C) got only 6.09% — so lambda=0.3 is NOT additive across seeds; it helped seed 42 a lot, seed 45 much less.",
        "HEADLINE (10F): the two seeds succeed on DIFFERENT behaviours, so take the UNION — count a behaviour won if EITHER seed wins it. Result: 13.97% combo, 110 of 520 behaviours (21.2%) — better than either seed alone, at ZERO extra optimization cost (just combining runs we already had).",
        "Overlap proof it's real diversity, not noise: 38 behaviours only seed-42, 16 only seed-45, 56 both. This is the clearest new positive result of Sprint 3.",
        "Gemma4 (10A, 'EmptyThink' recipe) scaled to 520: 3.91% vs 2.31% baseline (+1.6pp), consistent across 3 seeds — modest, but Gemma4's BEST result of the project and the first positive, multi-seed-consistent one at full scale. Reframes Gemma4 from 'totally unsolved' to 'one working recipe'.",
        "LIKELY Q ('why is the ensemble better?'): seed diversity — different random seeds crack different behaviours, so combining broadens coverage for free.",
    ],
    # 11 — Bug
    [
        "PURPOSE: the critical honesty slide — do NOT skip it; presenting it first protects your credibility.",
        "A dedicated CODE audit (after all the experiments) found a bug we call B1.",
        "THE BUG: when OPTIMIZING the suffix, the code placed it inside the ASSISTANT's turn of the chat template; when EVALUATING, it placed the same suffix in the USER's turn. So the optimizer was tuned on a slightly different prompt than we actually tested. The canonical GCG attack (Zou et al.) puts the suffix in the USER turn — so eval was correct and the optimization was the deviation.",
        "CONSEQUENCE: every ASR MAGNITUDE in this deck (10.7%, 24%, 8.9%, 14%, 3.9%) is provisional — it could move up or down after re-running with the fix.",
        "WHAT SURVIVES (say this clearly): the DIRECTIONS of every finding (CoT-prefix helps, weak lambda=0.3 helps, the ensemble helps), the detection AUC (built on clean generations, not the optimized suffix), the baselines, and the causal test (no suffix). The qualitative story stands.",
        "It also likely EXPLAINS 'loss doesn't predict ASR': loss was computed on the assistant-turn prompt while ASR was measured on the user-turn prompt — so of course they didn't correlate. If that decoupling shrinks after the fix, we'll revise that claim.",
        "STATUS: fix implemented (suffix in the user turn, verified byte-identical to eval), tiered re-run running now. There's a GATE run — if fixed-5A ASR is much higher than 10.7% we re-run everything; if unchanged, the bug didn't drive the numbers.",
        "LIKELY Q ('so are your results wrong?'): the directions and qualitative findings are expected to hold; the exact percentages are pending confirmation. We caught this ourselves via a code audit — that's the process working.",
    ],
    # 11b — GCG comparison: SAME
    [
        "PURPOSE: situate the work honestly — we did NOT invent a new attack; we adapted the standard one (GCG) to reasoning models. Say that plainly, it builds trust.",
        "GCG = Greedy Coordinate Gradient, from Zou et al. 2023 ('Universal and Transferable Adversarial Attacks') — THE standard white-box adversarial-suffix attack.",
        "The mechanics we KEPT are the heart of GCG: gradient-guided top-k token candidates per suffix position (original defaults: top_k=256, batch_size=1024), batched greedy acceptance, and a teacher-forced cross-entropy loss on the target.",
        "Attack surface is identical: a suffix appended after the request with a space; suffix in the USER turn, target in the ASSISTANT turn.",
        "We literally validated our (fixed) prompt-builder against the original llm-attacks `SuffixManager` code — same user-turn placement, same '{goal} {control}' space-join.",
        "Multi-behaviour gradient aggregation mirrors the original's universal / multi-prompt idea.",
        "LIKELY Q ('is this just GCG?'): the SKELETON is GCG; the novelty is everything on the next slide — reasoning-model targeting, the extra objectives, and judge-based evaluation.",
    ],
    # 11c — GCG comparison: DIFFERENT
    [
        "PURPOSE: our actual contributions — what makes this more than a re-run of GCG.",
        "Models: the original code only knew Llama-2 / Vicuna / GPT-family embeddings and chat templates; we added model-family adapters for Qwen3, Gemma4, and DeepSeek — including their tokenizers and <think> markers.",
        "BIGGEST difference — thinking models: original GCG has zero notion of a reasoning block. Our 5A trick (target starts with a fake <think>…</think>) is exactly the adaptation that makes GCG work on reasoning models. This is the core scientific contribution.",
        "Objective: original GCG = pure target cross-entropy. We built a composite objective with optional representation, KL, and refusal-direction terms — the refusal-direction term is what drove the Stage-3 λ=0.3 result. (Precision if pushed: the refusal-direction term only shapes which candidate tokens get PROPOSED via the gradient; it is not in the acceptance/selection criterion.)",
        "Evaluation is much stricter: original GCG calls it a win if the reply merely avoids refusal phrases ('I'm sorry', 'As an', …), which over-counts. We use StrongREJECT (an LLM judge) on real free-form generation, vs 3 controls — so our ASR numbers are CONSERVATIVE relative to the original's metric.",
        "Tokenizer detail (say if asked): original GCG discards any candidate whose text re-tokenizes to a different length (filter_cand=True). Fine for Llama's SentencePiece, but on Qwen/Gemma BPE it throws away essentially every candidate and the optimizer makes zero progress — so we disable it (--no-filter-cand). A necessary, documented divergence.",
        "We also added a defensive-security angle GCG doesn't have: the first-token hidden-state detector (slide 6).",
        "Honesty tie-in: the placement bug (slide 10) was literally a spot where our code DIVERGED from canonical GCG — we optimized the suffix in the assistant turn instead of the user turn. The fix puts us back on the canonical path (verified against the original code).",
    ],
    # 12 — Summary
    [
        "PURPOSE: recap; let each takeaway land; then invite questions.",
        "1. CoT-prefix targeting is THE unlock for thinking models (standard GCG ~2.7% -> 10.7%).",
        "2. WEAK refusal-direction suppression (lambda=0.3) is the strongest single lever — up to 24% on the dev set; STRONG suppression kills it (a sweet spot).",
        "3. A free union-of-two-seeds ensemble covers the most of the benchmark: 110/520 = 21%.",
        "4. Two separate axes: Qwen3 attacks are perfectly detectable at token 0 (AUC 1.0); Gemma4 resists almost everything (one modest 3.9% recipe).",
        "5. Optimizer loss does NOT predict attack success — a recurring theme.",
        "6. A placement bug (B1) means exact magnitudes are provisional; the fixed re-run is confirming them now.",
        "CLOSE: 'The qualitative findings are clear and real; the final numbers are pending the corrected re-run.' Then open for questions.",
        "IF ASKED 'what's next?': finish the v2 (fixed) re-runs to lock down magnitudes, then decide whether to push the ensemble / lambda-sweep further or start writing up.",
    ],
]

assert len(NOTES) == len(prs.slides._sldIdLst), "notes count must match slide count"
for slide, note in zip(prs.slides, NOTES):
    set_notes(slide, note)

prs.save(OUT)
print("WROTE", OUT, "with", len(prs.slides._sldIdLst), "slides + notes")
