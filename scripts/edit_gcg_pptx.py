#!/usr/bin/env python3
"""Apply audit corrections to the GCG Phases 4-7 deck.

Each slide's paragraphs are single-run (confirmed by inspection), so simple
run.text substring replacement preserves all existing formatting/theme.

Run under the base env (has python-pptx; does NOT need torch/sklearn):
  python3 scripts/edit_gcg_pptx.py

Reads:  docs/GCG_Phases4-7_Summary_2026-07-11.pptx  (untouched, left as-is)
Writes: docs/GCG_Phases4-7_Summary_FINAL_AUDITED.pptx  (new file)
"""
from pptx import Presentation

SRC = "docs/GCG_Phases4-7_Summary_2026-07-11.pptx"
DST = "docs/GCG_Phases4-7_Summary_FINAL_AUDITED.pptx"

# (substring_to_find, replacement_text) -- applied via str.replace on each run's
# full text. Order matters only in that later replacements operate on already-
# modified text, so keep each entry self-contained.
REPLACEMENTS = [
    # --- Slide 4: seed terminology definition (table cells, handled separately from row text) ---
    (
        "Training vs. unseen seeds",
        "Optimization vs. generation seeds",
    ),
    (
        "\"Training seeds\" (42/43/44) were used during suffix optimization for baseline comparisons; \"unseen seeds\" (100/200/300) are held out entirely, used only to test whether the suffix generalizes to fresh randomness",
        "\"Optimization seed\" controls which suffix GCG finds (varied 42/43/44/45 in 7B). \"Generation seed\" controls sampling for a fixed suffix: a dev panel {42,43,44} is used for all runs' baseline/comparison numbers (not necessarily matching that run's own optimization seed), and {100,200,300} are held out entirely to test generalization to fresh randomness",
    ),
    # --- Slide 8: Phase 6 detail -- 6A-Q target mislabel + softened framing ---
    (
        "Phase 6: Refusal-Direction Is Mutually Destructive",
        "Phase 6: Refusal-Direction Suppression Eliminates the CoT-Prefix Gain",
    ),
    (
        "6A-Q (Qwen3, CoT-prefix + refusal-direction loss added): 0% ASR, net-negative (−10.7pp vs. 5A alone)",
        "6A-Q (Qwen3, STANDARD target + refusal-direction loss added — not CoT-prefix): 0% ASR, net-negative (−2.7pp vs. task_only)",
    ),
    (
        "→ Conclusion: forcing the suffix to ALSO suppress the refusal direction pushes the optimizer toward unusual token combinations that themselves look suspicious to the model — the two objectives (predict compliant text AND suppress refusal) fight each other in token space",
        "→ Conclusion: under this tested setup (one layer=25, one lambda, one seed; the refusal-direction loss shapes the gradient's candidate proposals but is not part of the candidate-selection objective — see docs/GCG_REFUSAL_DIRECTION_AUDIT.md), adding refusal-direction suppression eliminated the CoT-prefix ASR gain, despite the optimizer successfully reducing the measured (single-behavior) refusal-direction projection. A broader claim of fundamental incompatibility would need a layer/lambda/seed sweep, not done here.",
    ),
    # --- Slide 10: 7A coverage phrasing ---
    (
        "FINAL unseen-seed result (5,849 rows, 493/520 behaviors reached, seeds 100/200/300): optimized_weighted = 8.92% ASR (131/1468) vs. neutral_control 3.83%  →  +5.09pp uplift — essentially the same rate as training seeds, confirming the suffix generalizes to fresh randomness too",
        "FINAL unseen-seed result (5,849 rows, 493/520 behaviors = 94.8% coverage, seeds 100/200/300): optimized_weighted = 8.92% ASR (131/1468) vs. neutral_control 3.83%  →  +5.09pp uplift over the 493 completed behaviors (missing 27 are the unfinished tail of specific shards, not a random sample — see docs/GCG_7A_BEHAVIOR_LEVEL_ANALYSIS.md; report as coverage-qualified, not an unbiased full-520 estimate)",
    ),
    # --- Slide 11: 7B seed table caption clarifying dev-panel terminology ---
    (
        "Same setup as 5A (500 optimization steps), only the random seed changed — this measures how much the attack's success depends on lucky/unlucky initialization",
        "Same setup as 5A (500 optimization steps), only the optimization seed changed. \"Training ASR\" in the table is measured on a fixed generation-seed dev panel {42,43,44} shared across all four rows — not each row's own optimization seed (seed 45's own generation seed was never sampled anywhere in this project; see docs/GCG_PHASE4_7_AUDIT_REPORT.md item 5)",
    ),
    # --- Slide 13: Gemma4 intrinsic-robustness overclaim ---
    (
        "→ Conclusion: format was NOT the barrier. Gemma4's resistance is intrinsic — attributed to a 58% stronger internal 'refusal direction' signal (0.498 vs. Qwen3's 0.315) and safety that is not reducible to a single interpretable direction or CoT-alignment issue",
        "→ Conclusion: the channel-token format-mismatch hypothesis is ruled out — Gemma4 resisted all 5 tested GCG configurations. It has a measurably stronger internal 'refusal direction' signal (0.498 vs. Qwen3's 0.315), an observed representational difference; whether that difference (or a more distributed safety mechanism) causally explains the robustness is a hypothesis for future work, not established by this test",
    ),
    # --- Slide 17 (Finding 3): title + closing line ---
    (
        "Finding 3: Refusal-Direction Suppression Is Mutually Destructive",
        "Finding 3: Refusal-Direction Suppression Eliminated the CoT-Prefix Gain (Tested Setup)",
    ),
    (
        "The two objectives — predict compliant text via cross-entropy, AND suppress the refusal direction — are fundamentally incompatible within the space of valid suffix tokens",
        "Under the one layer / one lambda / one optimization-seed setup tested here, the two objectives did not combine productively in valid suffix-token space; a broader incompatibility claim would need the layer/lambda/seed sweep flagged as future work in docs/GCG_REFUSAL_DIRECTION_AUDIT.md",
    ),
    # --- Slide 18 (Finding 4): title + body ---
    (
        "Finding 4: Gemma4 Is Intrinsically Robust (Not a Format Artifact)",
        "Finding 4: Gemma4 Resisted Every Tested GCG Variant (Format Hypothesis Ruled Out)",
    ),
    (
        "Gemma4's safety also appears to be distributed across multiple mechanisms rather than reducible to one 1-D direction or one CoT-alignment quirk — which is why even the most successful Qwen3 attack recipe, applied optimally, cannot move it off 0%",
        "Whether Gemma4's safety is distributed across multiple mechanisms rather than reducible to one 1-D direction is a hypothesis suggested by these results, not directly tested (would need activation-patching analogous to the prior paper's Qwen3 ablation, applied to Gemma4)",
    ),
    # --- Slide 22: Implications point 3 ---
    (
        "3. Gemma4's multi-layered safety appears meaningfully harder to attack via GCG than Qwen3's: a 58% stronger refusal-direction signal plus safety that isn't reducible to a single interpretable direction. This suggests architectural or training choices that produce more distributed, harder-to-locate-and-suppress safety representations are a promising direction for hardening other models",
        "3. Gemma4 was meaningfully harder to attack via GCG than Qwen3 across every configuration tried, alongside a 58% stronger measured refusal-direction signal. Whether more \"distributed\" safety representations are the causal reason (vs. other architectural/training differences) is a hypothesis worth testing directly (e.g. via activation-patching) before it informs hardening recommendations for other models",
    ),
    # --- Slide 23: status table wording ---
    (
        "Refusal-direction combination confirmed mutually destructive",
        "Refusal-direction + GCG eliminated the CoT-prefix ASR gain under the tested setup (one layer/lambda/seed) — see docs/GCG_REFUSAL_DIRECTION_AUDIT.md",
    ),
    (
        "Confirms intrinsic robustness, not format artifact",
        "Rules out the channel-token format-mismatch hypothesis; Gemma4 resisted all 5 tested configurations",
    ),
    # --- Slide 24: conclusions bullets 3 and 4 ---
    (
        "3. Combining GCG with refusal-direction suppression backfires — 0% net-negative, vs. 91% for direct activation ablation in prior work",
        "3. Under the tested setup, combining GCG with refusal-direction suppression eliminated the CoT-prefix ASR gain (0% net-negative vs. 91% for direct activation ablation in prior work, a different, white-box intervention) — not established as a fundamental incompatibility beyond this configuration",
    ),
    (
        "4. Gemma4 resisted every attack variant tried (0% ASR universally, 5 distinct configurations) — confirmed intrinsic, not a CoT-format artifact",
        "4. Gemma4 resisted every attack variant tried (0% ASR, 5 distinct configurations) — the CoT-format-mismatch hypothesis is ruled out; deeper causal explanation remains a hypothesis",
    ),
    # --- Slide 14: title, add asterisk footnote pointer for the 6A/Qwen3 baseline fix above ---
    (
        "Master Results Table — All Experiments, Phases 4–7",
        "Master Results Table — All Experiments, Phases 4–7 (*6A/Qwen3: vs. task_only, since it uses the standard target, not CoT-prefix)",
    ),
    # --- Slide 25: fix stale shard/wave count (written before unseeded pass 3 was run) ---
    (
        "‒  scripts/split_unseeded_shards.py / merge_unseeded_shards.py — same pattern for the 3-unseen-seed evaluation (10 shards across 2 waves)",
        "‒  scripts/split_unseeded_shards.py / merge_unseeded_shards.py — same pattern for the 3-unseen-seed evaluation (15 shards across 3 waves, matching slide 10 — this bullet was stale from before the 3rd wave ran, corrected 2026-07-13)",
    ),
    # --- Slide 25: appendix, add audit docs pointer ---
    (
        "‒  docs/GCG_ABLATION_PIPELINE_LOG.md — detailed Phase 4–6 job-by-job log",
        "‒  docs/GCG_ABLATION_PIPELINE_LOG.md — detailed Phase 4–6 job-by-job log; docs/GCG_PHASE4_7_AUDIT_REPORT.md + docs/GCG_PHASE4_7_SOURCE_OF_TRUTH.md (table csv in outputs/stage_gcg_full/) + 4 companion audit docs (2026-07-13 audit, source of truth for every number in this deck)",
    ),
]


def _text_frames_in_shape(shape):
    if shape.has_text_frame:
        yield shape.text_frame
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame


def fix_slide14_master_table_6a_baseline(prs):
    """Slide 14's master table 6A/Qwen3 row still showed the pre-correction '-10.7pp'
    baseline delta (comparing to 5A's CoT-prefix ASR) after slide 8 was corrected to
    identify 6A-Q as using the STANDARD target -- its correct comparison is vs. task_only
    at -2.7pp. Row-scoped fix (not a generic text replace) since 6C's row legitimately
    keeps '-10.7pp' (6C genuinely uses the CoT-prefix target, so vs-5A is the right
    comparison there) -- found by independent post-hoc verification, 2026-07-13."""
    slide14 = prs.slides[13]
    fixed = 0
    for shape in slide14.shapes:
        if not shape.has_table:
            continue
        for row in shape.table.rows:
            cells = list(row.cells)
            cells_text = [c.text for c in cells]
            if cells_text[:2] == ["6A", "Qwen3"] and cells_text[-1].strip() == "−10.7pp":
                last_cell = cells[-1]
                for para in last_cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip() == "−10.7pp":
                            run.text = "−2.7pp*"
                            fixed += 1
    return fixed


def apply_replacements(prs):
    counts = {i: 0 for i in range(len(REPLACEMENTS))}
    for slide in prs.slides:
        for shape in slide.shapes:
            for tf in _text_frames_in_shape(shape):
                for para in tf.paragraphs:
                    for run in para.runs:
                        for i, (old, new) in enumerate(REPLACEMENTS):
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                counts[i] += 1
    return counts


def main():
    prs = Presentation(SRC)
    counts = apply_replacements(prs)
    n_fixed_14 = fix_slide14_master_table_6a_baseline(prs)
    prs.save(DST)
    print(f"Saved {DST}")
    for i, (old, _new) in enumerate(REPLACEMENTS):
        status = "OK" if counts[i] > 0 else "*** NOT FOUND ***"
        print(f"[{counts[i]}] {status}  {old[:70]!r}")
    print(f"[{n_fixed_14}] slide14_6A_baseline_cell_fix {'OK' if n_fixed_14 else '*** NOT FOUND ***'}")


if __name__ == "__main__":
    main()
