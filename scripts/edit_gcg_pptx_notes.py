#!/usr/bin/env python3
"""Correct stale/overclaiming speaker notes on the audited GCG deck.

The original deck already had substantial speaker notes (600-1500 chars/slide),
written before this audit. Several of them state now-corrected claims as
established fact (most seriously: slide 10's notes call the 7A unseeded 8.92%
estimate "unbiased" -- the exact claim this audit determined was NOT supported,
since the 27 missing behaviors are the non-random tail of interrupted shards).
This script applies targeted corrections to those notes, operating on
docs/GCG_Phases4-7_Summary_FINAL_AUDITED.pptx IN PLACE (it is already the
corrected/working copy, distinct from the untouched original
docs/GCG_Phases4-7_Summary_2026-07-11.pptx).

Run order matters: this must run AFTER scripts/edit_gcg_pptx.py (which regenerates
DECK from the untouched original each time it runs, and would silently wipe out
these notes fixes if rerun afterward). Run: python3 scripts/edit_gcg_pptx.py
&& python3 scripts/edit_gcg_pptx_notes.py
"""
from pptx import Presentation

DECK = "docs/GCG_Phases4-7_Summary_FINAL_AUDITED.pptx"

# (substring, replacement) applied to notes_text_frame runs.
NOTES_REPLACEMENTS = [
    # --- Slide 8 (Phase 6 detail): hedge the "other safety pathways" claim, note 6A-Q target fix ---
    (
        "The paper's method has essentially unconstrained access to internal activations; GCG is constrained to whatever effect real tokens can have on those same activations — and apparently, tokens weird enough to suppress the refusal direction are also weird enough to trigger the model's other, non-refusal-direction safety pathways.",
        "The paper's method has essentially unconstrained access to internal activations; GCG is constrained to whatever effect real tokens can have on those same activations. One candidate explanation is that tokens weird enough to suppress the refusal direction are also weird enough to trigger other safety pathways -- but this project didn't isolate that mechanism directly, so treat it as a hypothesis, not a confirmed explanation. Also note: 6A-Q (this slide) uses the STANDARD target, not the CoT-prefix target -- only 6C (slide 14) combines both. And under the hood, the refusal-direction loss only shapes which candidates get proposed via the gradient; it's not part of the actual candidate-selection criterion each step (see the 2026-07-13 audit's refusal-direction doc for the full trace).",
    ),
    # --- Slide 10 (7A results): the single most important notes fix -- "unbiased" claim corrected ---
    (
        "The unseen-seed evaluation reached 493 of 520 behaviors (94.8%) — 27 behaviors were not reached before the wall-time limit across all three passes, but assignment of behaviors to shards was approximately random, so the 8.92% estimate is considered unbiased.",
        "The unseen-seed evaluation reached 493 of 520 behaviors (94.8%) — 27 behaviors were not reached because 4 of the parallel shard jobs ran out of wall-time before finishing their assigned range; a 2026-07-13 audit found this missingness is NOT random (the 27 missing behaviors are specifically the trailing tail of those shards, skewing toward higher AdvBench row indices, and mildly over-representing the misinformation category), so 8.92% should be presented as \"over the 493 completed behaviors (94.8% coverage),\" not as an unbiased estimate of the full 520.",
    ),
    # --- Slide 17 (Finding 3): hedge "more than one way of encoding" claim, scope to tested config ---
    (
        "A technique that achieves 91% success when directly patching activations achieves 0% when the same underlying idea has to be expressed indirectly through token choice — the model has more than one way of encoding 'this looks unsafe', and pushing on one dimension (the refusal direction) via token choice triggers others.",
        "A technique that achieves 91% success when directly patching activations achieves 0% when the same underlying idea has to be expressed indirectly through token choice, under the one layer/lambda/optimization-seed combination tested here. A plausible explanation is that the model has more than one way of encoding 'this looks unsafe' and pushing on one dimension (the refusal direction) via token choice triggers others -- but this project didn't test enough layers/lambdas/seeds to confirm that generally, so present it as a hypothesis consistent with the result, not a proven mechanism.",
    ),
    # --- Slide 18 (Finding 4): soften "best-supported conclusion" to hypothesis ---
    (
        "The 7C experiment (slide 13) was the decisive test — it deliberately removed the one remaining plausible confound (CoT-format misalignment, the same bug that explained Qwen3's initial weakness) and still found 0%, which rules out the simplest alternative explanation and leaves 'genuinely more robust safety training' as the best-supported conclusion.",
        "The 7C experiment (slide 13) was the decisive test — it deliberately removed the one remaining plausible confound (CoT-format misalignment, the same bug that explained Qwen3's initial weakness) and still found 0%, which rules out that specific alternative explanation. What's left is that Gemma4 resisted every configuration tried; whether that reflects deeper/more distributed safety training, as opposed to some other architectural or training difference, remains a hypothesis this project didn't test directly (e.g. via activation-patching analogous to the paper's Qwen3 ablation).",
    ),
    # --- Slide 22 (Implications): hedge "safety not concentrated in one direction" ---
    (
        "Recommendation #3 is more of a research direction for model builders: Gemma4's resistance profile (harder to attack, safety not concentrated in one direction) is worth understanding well enough to potentially replicate in other models' safety training.",
        "Recommendation #3 is more of a research direction for model builders: Gemma4 was harder to attack across every configuration tried, alongside a measurably stronger refusal-direction signal -- whether its safety is genuinely \"not concentrated in one direction\" is a hypothesis worth testing directly (not yet confirmed here) before treating it as a template to replicate in other models' safety training.",
    ),
    # --- Slide 23 (Current Status): stale shard/wave count, same bug class already fixed on slide 25's body text ---
    (
        "7A ended up being the longest-running experiment because of its sheer size (520 behaviors, and later a second unseen-seed pass on top), so it was split into parallel SLURM jobs on the cluster — first 6 shards for the main evaluation, then 10 more shards across two waves for the unseen-seed evaluation — using helper scripts",
        "7A ended up being the longest-running experiment because of its sheer size (520 behaviors, and later a second unseen-seed pass on top), so it was split into parallel SLURM jobs on the cluster — first 6 shards for the main evaluation, then 15 more shards across three waves for the unseen-seed evaluation (a 3rd wave was needed to close the remaining gap, though 27 behaviors still weren't reached by wall-time — see slide 10) — using helper scripts",
    ),
    # --- Slide 25 (Appendix): point to the new audit docs ---
    (
        "This appendix is a pointer reference, not meant to be read top-to-bottom — it's here so anyone who wants to dig into the raw underlying data after seeing this summary knows exactly which files and job IDs to look for.",
        "This appendix is a pointer reference, not meant to be read top-to-bottom — it's here so anyone who wants to dig into the raw underlying data after seeing this summary knows exactly which files and job IDs to look for. A 2026-07-13 audit pass produced six additional documents (GCG_PHASE4_7_AUDIT_REPORT.md, GCG_PHASE4_7_SOURCE_OF_TRUTH.{csv,md}, GCG_DETECTOR_ROBUSTNESS_AUDIT.md, GCG_REFUSAL_DIRECTION_AUDIT.md, GCG_COT_PREFIX_MECHANISM_ANALYSIS.md, GCG_7A_BEHAVIOR_LEVEL_ANALYSIS.md) that trace every number in this deck back to raw artifacts and correct several claims that were previously overstated -- start there for anything that seems too strong or too clean.",
    ),
]


def apply_notes_replacements(prs):
    counts = {i: 0 for i in range(len(NOTES_REPLACEMENTS))}
    for slide in prs.slides:
        if not slide.has_notes_slide:
            continue
        tf = slide.notes_slide.notes_text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                for i, (old, new) in enumerate(NOTES_REPLACEMENTS):
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        counts[i] += 1
    return counts


def main():
    prs = Presentation(DECK)
    counts = apply_notes_replacements(prs)
    prs.save(DECK)
    print(f"Saved {DECK} (in place)")
    for i, (old, _new) in enumerate(NOTES_REPLACEMENTS):
        status = "OK" if counts[i] > 0 else "*** NOT FOUND ***"
        print(f"[{counts[i]}] {status}  {old[:70]!r}")


if __name__ == "__main__":
    main()
