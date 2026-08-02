"""Shared style system for the research-update deck.

One font (Arial), one palette, one grid. No page numbers, no eyebrows, no footers.
"""
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
import copy

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x15, 0x3A, 0x5C)   # titles, table header fill, rule
INK = RGBColor(0x1C, 0x22, 0x2E)   # body text
MUTED = RGBColor(0x6A, 0x72, 0x7D)   # secondary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_A = "F3F5F7"
ROW_B = "E9EDF1"
HDR = "153A5C"

FONT = "Arial"

# ---------------------------------------------------------------- grid (EMU)
SLIDE_W = 12192000
SLIDE_H = 6858000
L = 566928                 # left margin
CW = 11057839              # content width
RULE_Y = 1600200           # horizontal rule under the title
TITLE_TOP = 594360
TITLE_H = RULE_Y - TITLE_TOP - 137160
BODY_TOP = 1828800
BOTTOM = 6400800           # content must not cross this


def new_deck():
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    # drop the default slide layouts we do not use; keep Blank (index 6)
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _style_run(r, size, bold=False, color=INK, italic=False):
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    # force the East-Asian / complex-script faces too so nothing falls back
    rPr = r._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        el = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag)
        if el is None:
            el = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag)
        el.set("typeface", FONT)


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, bold=False, color=INK, space_after=8, first=False,
         line=1.15, align=PP_ALIGN.LEFT, italic=False):
    """Add a paragraph. `text` may be a list of (chunk, bold) tuples for mixed runs."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    chunks = text if isinstance(text, list) else [(text, bold)]
    for chunk, b in chunks:
        r = p.add_run()
        r.text = chunk
        _style_run(r, size, bold=b, color=color, italic=italic)
    return p


def title(slide, text, size=23):
    tb, tf = textbox(slide, L, TITLE_TOP, CW, TITLE_H, anchor=MSO_ANCHOR.BOTTOM)
    para(tf, text, size, bold=True, color=NAVY, space_after=0, first=True, line=1.05)
    rule(slide)
    return tb


def rule(slide, y=RULE_Y):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(L), Emu(y), Emu(L + CW), Emu(y))
    c.line.color.rgb = NAVY
    c.line.width = Emu(22860)
    return c


def body(slide, lines, size=15.5, top=BODY_TOP, width=CW, space_after=10, line=1.2):
    """lines: list of strings or list-of-(chunk,bold)."""
    tb, tf = textbox(slide, L, top, width, 300000)
    for i, ln in enumerate(lines):
        para(tf, ln, size, space_after=space_after, first=(i == 0), line=line)
    return tb


def _split_bold(text, terms):
    """Split text into (chunk, bold) runs, bolding every occurrence of any term."""
    import re
    if not terms:
        return [(text, False)]
    pat = re.compile("(" + "|".join(re.escape(t) for t in sorted(terms, key=len, reverse=True)) + ")")
    return [(part, bool(pat.fullmatch(part))) for part in pat.split(text) if part]


def _set_cell(cell, text, size, bold=False, color=INK, fill=None, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.MIDDLE, line=1.05, bold_terms=None):
    cell.margin_left = Emu(88900)
    cell.margin_right = Emu(88900)
    cell.margin_top = Emu(38100)
    cell.margin_bottom = Emu(38100)
    cell.vertical_anchor = anchor
    tf = cell.text_frame
    tf.word_wrap = True
    parts = text.split("\n")
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(0)
        for chunk, is_b in _split_bold(part, bold_terms):
            r = p.add_run()
            r.text = chunk
            _style_run(r, size, bold=(bold or is_b), color=color)
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        cell.fill.background()


def _kill_theme_banding(tbl):
    """python-pptx applies a banded medium style by default; neutralise it."""
    tblPr = tbl._tbl.find("{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr")
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    for el in tblPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"):
        tblPr.remove(el)
    sid = etree.SubElement(tblPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId")
    sid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # "No Style, No Grid"


EMU_PT = 12700
_LIB_REG = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
_LIB_BOLD = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
_fcache = {}


def _lib(size_pt, bold):
    """Liberation Sans is metrically identical to Arial, so its extents are Arial's."""
    from PIL import ImageFont
    key = (round(size_pt * 4), bold)
    if key not in _fcache:
        _fcache[key] = ImageFont.truetype(_LIB_BOLD if bold else _LIB_REG, int(round(size_pt * 4)))
    return _fcache[key]


def _n_lines(text, size_pt, bold, avail_pt):
    total = 0
    f = _lib(size_pt, bold)
    for hard in text.split("\n"):
        line, n = "", 1
        for w in hard.split(" "):
            trial = (line + " " + w).strip()
            if f.getlength(trial) / 4.0 <= avail_pt or not line:
                line = trial
            else:
                n += 1
                line = w
        total += n
    return total


def autofit_rows(rows, col_widths, size, hdr_size, header, label_col, pad_pt=7.0,
                 margin=88900):
    """Row heights that exactly fit the wrapped text — no PowerPoint auto-grow surprises."""
    heights = []
    for ri, row in enumerate(rows):
        is_hdr = header and ri == 0
        need = 0.0
        for ci, val in enumerate(row):
            sz = hdr_size if is_hdr else size
            bold = is_hdr or (label_col and ci == 0)
            avail = (col_widths[ci] - 2 * margin) / EMU_PT
            need = max(need, _n_lines(str(val), sz, bold, avail) * sz * 1.24)
        heights.append(int(round((need + pad_pt) * EMU_PT)))
    return heights


def table(slide, rows, col_widths, top, header=True, size=10.5, hdr_size=11.5,
          label_col=True, row_h=None, align_first=PP_ALIGN.LEFT, num_align=PP_ALIGN.LEFT,
          left=L, autofit=False, min_row_h=0, fill_to=None, bold_terms=None,
          anchor=MSO_ANCHOR.MIDDLE):
    n_rows, n_cols = len(rows), len(rows[0])
    if autofit:
        heights = autofit_rows(rows, col_widths, size, hdr_size, header, label_col)
        heights = [max(h, min_row_h) for h in heights]
        if fill_to:
            slack = fill_to - top - sum(heights)
            if slack > 0:
                per = slack // len(heights)
                heights = [h + per for h in heights]
        row_h = None
    total_w = sum(col_widths)
    shp = slide.shapes.add_table(n_rows, n_cols, Emu(left), Emu(top), Emu(total_w),
                                 Emu(row_h * n_rows if row_h else 400000 * n_rows))
    tbl = shp.table
    _kill_theme_banding(tbl)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Emu(w)
    if row_h:
        for r in tbl.rows:
            r.height = Emu(row_h)
    elif autofit:
        for r, h in zip(tbl.rows, heights):
            r.height = Emu(h)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if header and ri == 0:
                _set_cell(cell, val, hdr_size, bold=True, color=WHITE, fill=HDR,
                          align=align_first if ci == 0 else num_align)
            else:
                fill = ROW_A if (ri % 2 == (1 if header else 0)) else ROW_B
                is_label = label_col and ci == 0
                _set_cell(cell, val, size, bold=is_label,
                          color=NAVY if is_label else INK, fill=fill,
                          align=align_first if ci == 0 else num_align,
                          bold_terms=None if is_label else bold_terms, anchor=anchor)
    return shp


def notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text
