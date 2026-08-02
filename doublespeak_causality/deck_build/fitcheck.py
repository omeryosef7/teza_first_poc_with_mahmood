#!/usr/bin/env python
"""Overflow checker: Liberation Sans is metrically identical to Arial, so text
extents measured with it are the extents PowerPoint will use."""
import sys
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

EMU_PT = 12700
REG = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
BOLD = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
_cache = {}


def font(size_pt, bold):
    key = (round(size_pt * 4), bold)
    if key not in _cache:
        _cache[key] = ImageFont.truetype(BOLD if bold else REG, int(round(size_pt * 4)))
    return _cache[key]


def width_pt(text, size_pt, bold):
    return font(size_pt, bold).getlength(text) / 4.0


def wrap_lines(text, size_pt, bold, avail_pt):
    total = 0
    for hard in text.split("\n"):
        words, line, n = hard.split(" "), "", 1
        for w in words:
            trial = (line + " " + w).strip()
            if width_pt(trial, size_pt, bold) <= avail_pt or not line:
                line = trial
            else:
                n += 1
                line = w
        total += n
    return total


def check(path):
    prs = Presentation(path)
    SH = prs.slide_height
    bad = []
    for i, s in enumerate(prs.slides, 1):
        spans = []          # (top, bottom, label) of every rendered block on this slide
        for sh in s.shapes:
            if sh.has_table:
                t = sh.table
                widths = [c.width for c in t.columns]
                tbl_top = sh.top
                total_h = 0
                for ri, row in enumerate(t.rows):
                    need = 0
                    for ci in range(len(widths)):
                        cell = t.cell(ri, ci)
                        avail = (widths[ci] - cell.margin_left - cell.margin_right) / EMU_PT
                        txt = cell.text_frame.text
                        runs = [r for p in cell.text_frame.paragraphs for r in p.runs]
                        if not runs:
                            continue
                        sz = runs[0].font.size.pt
                        bold = bool(runs[0].font.bold)
                        lines = wrap_lines(txt, sz, bold, avail)
                        h = lines * sz * 1.22 + (cell.margin_top + cell.margin_bottom) / EMU_PT
                        need = max(need, h)
                    have = row.height / EMU_PT
                    total_h += max(need, have)
                    if need > have + 0.5:
                        bad.append(f"slide {i} table row {ri}: needs {need:.1f}pt, row height {have:.1f}pt")
                bottom = tbl_top + total_h * EMU_PT
                spans.append((tbl_top, bottom, "table"))
                flag = "  <== OFF-SLIDE" if bottom > SH - 114300 else ""
                print(f"slide {i}: table top={tbl_top/914400:.2f}in  height={total_h/72:.2f}in  "
                      f"bottom={bottom/914400:.2f}in{flag}")
                if bottom > SH - 114300:
                    bad.append(f"slide {i}: table bottom {bottom/914400:.2f}in > 7.38in")
            elif sh.has_text_frame and sh.text_frame.text.strip():
                tf = sh.text_frame
                avail = (sh.width - tf.margin_left - tf.margin_right) / EMU_PT
                y = 0
                for p in tf.paragraphs:
                    if not p.runs:
                        continue
                    sz = p.runs[0].font.size.pt
                    bold = bool(p.runs[0].font.bold)
                    txt = "".join(r.text for r in p.runs)
                    ls = p.line_spacing or 1.0
                    y += wrap_lines(txt, sz, bold, avail) * sz * 1.2 * ls
                    y += (p.space_after.pt if p.space_after else 0)
                bottom = sh.top + y * EMU_PT
                is_title = str(tf.vertical_anchor) == "BOTTOM (5)"
                if not is_title:
                    spans.append((sh.top, bottom, f"text {sh.text_frame.text[:20]!r}"))
                print(f"slide {i}: text '{sh.text_frame.text[:34]!r}...' top={sh.top/914400:.2f} "
                      f"needs={y/72:.2f}in bottom={bottom/914400:.2f}in")
                if bottom > prs.slide_height - 114300:
                    bad.append(f"slide {i}: textbox bottom {bottom/914400:.2f}in off slide")
            elif sh.shape_type == 13:
                spans.append((sh.top, sh.top + sh.height, "picture"))
        spans.sort()
        for (t1, b1, l1), (t2, b2, l2) in zip(spans, spans[1:]):
            if t2 < b1 - 9144:      # allow 0.01in slack
                bad.append(f"slide {i}: OVERLAP — {l1} ends {b1/914400:.2f}in but "
                           f"{l2} starts {t2/914400:.2f}in")
    print("\n=== PROBLEMS ===" if bad else "\n=== NO OVERFLOW ===")
    for b in bad:
        print(" -", b)
    return bad


if __name__ == "__main__":
    check(sys.argv[1])
